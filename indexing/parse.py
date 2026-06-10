"""포맷별 텍스트 추출 (PRD §10.2).

- PDF : pymupdf 페이지별 텍스트
- PPTX: python-pptx 슬라이드 본문 + 발표자 노트 (슬라이드 = 페이지)
- HWPX: LibreOffice headless 로 PDF 변환 후 처리 권장 (안내 메시지)
"""
from __future__ import annotations

import os
from dataclasses import dataclass


@dataclass
class Page:
    page_num: int | None
    section_title: str | None
    text: str


def parse_pdf(path: str) -> list:
    import fitz  # pymupdf

    doc = fitz.open(path)
    pages = []
    try:
        for i in range(doc.page_count):
            page = doc.load_page(i)
            text = page.get_text("text").strip()
            if text:
                pages.append(Page(page_num=i + 1, section_title=None, text=text))
    finally:
        doc.close()
    return pages


def parse_pptx(path: str) -> list:
    from pptx import Presentation

    prs = Presentation(path)
    pages = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
        # 발표자 노트
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append("[노트] " + notes)
        text = "\n".join(parts).strip()
        if text:
            pages.append(Page(page_num=i + 1, section_title=None, text=text))
    return pages


def parse_file(path: str) -> list:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return parse_pdf(path)
    if ext == ".pptx":
        return parse_pptx(path)
    if ext in (".hwpx", ".hwp"):
        raise NotImplementedError(
            "HWPX/HWP는 LibreOffice headless 로 PDF 변환 후 인덱싱하세요:\n"
            "  soffice --headless --convert-to pdf --outdir docs/ <파일>"
        )
    raise ValueError("지원하지 않는 형식: " + ext)
